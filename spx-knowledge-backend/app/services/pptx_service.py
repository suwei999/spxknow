"""
PowerPoint (PPTX) 文档解析服务
基于 python-pptx，输出结构与 DOCX/PDF 解析保持一致
"""

import os
from typing import Any, Dict, List, Optional
from io import BytesIO

from sqlalchemy.orm import Session
from app.core.logging import logger


class PptxService:
    """
    PowerPoint 文档解析服务，输出结构与 DOCX/PDF 解析保持一致。
    输出结构：
      - text_content: str
      - ordered_elements: List[Dict]
      - filtered_elements_light: List[Dict]
      - text_element_index_map: List[Dict]
      - tables: List[Dict]
      - images: List[Dict]
      - images_binary: List[Dict]
      - metadata: Dict
    """

    def __init__(self, db: Session):
        self.db = db

    def parse_document(self, file_path: str) -> Dict[str, Any]:
        """
        解析 PPTX 文档
        
        Args:
            file_path: PPTX 文件路径
            
        Returns:
            返回结构（与 DocxService 保持一致）
        """
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"文件不存在: {file_path}")
        
        try:
            from pptx import Presentation
        except Exception as e:
            raise RuntimeError("缺少依赖 python-pptx，请先安装: pip install python-pptx") from e

        logger.info(f"[PPTX] 开始解析: {file_path}")
        
        # 加载演示文稿
        try:
            prs = Presentation(file_path)
        except Exception as e:
            logger.error(f"[PPTX] 加载文件失败: {e}")
            raise RuntimeError(f"无法打开 PPTX 文件: {e}") from e

        ordered_elements: List[Dict[str, Any]] = []
        text_element_index_map: List[Dict[str, Any]] = []
        filtered_elements_light: List[Dict[str, Any]] = []
        text_content_parts: List[str] = []
        images_payload: List[Dict[str, Any]] = []
        images_binary: List[Dict[str, Any]] = []
        tables: List[Dict[str, Any]] = []

        element_index = 0
        doc_order = 0
        
        # 收集幻灯片信息
        slides_info: List[Dict[str, Any]] = []
        layout_types: set = set()
        has_notes = False
        table_count = 0
        image_count = 0
        has_smartart = False
        
        # 获取演示文稿尺寸
        try:
            slide_width = prs.slide_width
            slide_height = prs.slide_height
            # 转换为英寸（EMU 单位，1 英寸 = 914400 EMU）
            width_inches = slide_width / 914400.0 if slide_width else 10.0
            height_inches = slide_height / 914400.0 if slide_height else 7.5
            presentation_size = {"width": width_inches, "height": height_inches}
        except Exception:
            presentation_size = {"width": 10.0, "height": 7.5}

        # 遍历所有幻灯片
        for slide_idx, slide in enumerate(prs.slides):
            slide_number = slide_idx + 1
            
            try:
                # 识别幻灯片布局
                layout_name = self._identify_slide_layout(slide)
                layout_types.add(layout_name)
                
                # 构建幻灯片标题
                slide_title = self._build_slide_title(slide)
                
                # 提取幻灯片内容
                slide_content = self._extract_slide_content(slide, slide_number)
                
                # 处理文本内容（先处理内容，再收集信息，确保一致性）
                slide_text = slide_content.get("text", "")
                
                # 如果有标题，先添加标题（即使没有正文也要添加标题）
                if slide_title:
                    title_text = f"幻灯片 {slide_number}: {slide_title}\n\n"
                    element_index += 1
                    ordered_elements.append({
                        "type": "text",
                        "text": title_text.strip(),
                        "element_index": element_index,
                        "doc_order": doc_order,
                        "category": "Title",
                        "style": "",
                        "length": len(title_text),
                        "slide_number": slide_number,
                        "slide_title": slide_title,
                        "slide_layout": layout_name,
                    })
                    text_content_parts.append(title_text)
                    filtered_elements_light.append({
                        "category": "Title",
                        "text": title_text.strip(),
                        "element_index": element_index,
                        "doc_order": doc_order,
                    })
                    text_element_index_map.append({
                        "element_index": element_index,
                        "element_type": "Title",
                        "doc_order": doc_order,
                        "page_number": None,
                        "coordinates": None,
                    })
                    doc_order += 1
                
                # 添加正文内容（如果有）
                if slide_text and slide_text.strip():
                    element_index += 1
                    ordered_elements.append({
                        "type": "text",
                        "text": slide_text,
                        "element_index": element_index,
                        "doc_order": doc_order,
                        "category": "NarrativeText",
                        "style": "",
                        "length": len(slide_text),
                        "slide_number": slide_number,
                        "slide_title": slide_title,
                        "slide_layout": layout_name,
                    })
                    text_content_parts.append(slide_text)
                    filtered_elements_light.append({
                        "category": "NarrativeText",
                        "text": slide_text,
                        "element_index": element_index,
                        "doc_order": doc_order,
                    })
                    text_element_index_map.append({
                        "element_index": element_index,
                        "element_type": "NarrativeText",
                        "doc_order": doc_order,
                        "page_number": None,
                        "coordinates": None,
                    })
                    doc_order += 1
                
                # 处理表格
                slide_tables = slide_content.get("tables", [])
                for table_data in slide_tables:
                    table_count += 1
                    element_index += 1
                    
                    # 添加到 tables 列表
                    table_item = {
                        "element_index": element_index,
                        "doc_order": doc_order,
                        "table_data": table_data["table_data"],
                        "table_text": table_data["table_text"],
                        "slide_number": slide_number,
                        "page_number": None,
                    }
                    tables.append(table_item)
                    
                    # 添加到 ordered_elements
                    ordered_elements.append({
                        "type": "table",
                        "element_index": element_index,
                        "doc_order": doc_order,
                        "slide_number": slide_number,
                    })
                    
                    # 添加到 filtered_elements_light
                    filtered_elements_light.append({
                        "category": "Table",
                        "text": table_data["table_text"],
                        "element_index": element_index,
                        "doc_order": doc_order,
                    })
                    
                    doc_order += 1
                
                # 处理图片
                slide_images = slide_content.get("images", [])
                for image_data in slide_images:
                    image_count += 1
                    element_index += 1
                    
                    # 添加到 images 列表（必须包含 data 和 bytes）
                    image_item = {
                        "data": image_data["data"],
                        "bytes": image_data["data"],  # 与 data 相同
                        "element_index": element_index,
                        "doc_order": doc_order,
                        "page_number": None,
                        "slide_number": slide_number,
                        "image_ext": image_data.get("image_ext", ".png"),
                        "width": image_data.get("width"),
                        "height": image_data.get("height"),
                    }
                    images_payload.append(image_item)
                    images_binary.append({
                        "binary": image_data["data"],
                        "element_index": element_index,
                        "page_number": None,
                        "doc_order": doc_order,
                    })
                    
                    # 添加到 ordered_elements
                    ordered_elements.append({
                        "type": "image",
                        "element_index": element_index,
                        "doc_order": doc_order,
                        "slide_number": slide_number,
                        "image_ext": image_data.get("image_ext", ".png"),
                    })
                    
                    doc_order += 1
                
                # 处理备注
                notes_text = slide_content.get("notes", "")
                if notes_text:
                    has_notes = True
                    # 备注作为单独的文本元素
                    element_index += 1
                    notes_element = {
                        "type": "text",
                        "text": f"备注：{notes_text}",
                        "element_index": element_index,
                        "doc_order": doc_order,
                        "category": "NarrativeText",
                        "style": "",
                        "length": len(notes_text),
                        "slide_number": slide_number,
                        "slide_title": slide_title,
                        "slide_layout": layout_name,
                        "is_notes": True,
                    }
                    ordered_elements.append(notes_element)
                    text_content_parts.append(f"备注：{notes_text}")
                    filtered_elements_light.append({
                        "category": "NarrativeText",
                        "text": f"备注：{notes_text}",
                        "element_index": element_index,
                        "doc_order": doc_order,
                    })
                    text_element_index_map.append({
                        "element_index": element_index,
                        "element_type": "NarrativeText",
                        "doc_order": doc_order,
                        "page_number": None,
                        "coordinates": None,
                    })
                    doc_order += 1
                
                # 收集幻灯片信息（在处理完所有内容后）
                slide_info = {
                    "number": slide_number,
                    "title": slide_title or "",
                    "layout": layout_name,
                    "has_table": len(slide_content.get("tables", [])) > 0,
                    "has_images": len(slide_content.get("images", [])) > 0,
                }
                slides_info.append(slide_info)
                    
            except Exception as e:
                logger.warning(f"[PPTX] 处理幻灯片 {slide_number} 失败: {e}", exc_info=True)
                # 即使处理失败，也记录幻灯片信息（但标记为失败）
                slide_info = {
                    "number": slide_number,
                    "title": "",
                    "layout": "Unknown",
                    "has_table": False,
                    "has_images": False,
                    "error": str(e)[:200],  # 记录错误信息（截断）
                }
                slides_info.append(slide_info)
                # 继续处理下一张幻灯片
                continue

        # 构建元数据
        metadata = {
            "element_count": len(ordered_elements),
            "images_count": len(images_payload),
            "slide_count": len(prs.slides),
            "layout_types": sorted(list(layout_types)),
            "has_notes": has_notes,
            "table_count": table_count,
            "image_count": image_count,
            "has_smartart": has_smartart,
            "presentation_size": presentation_size,
            "slides": slides_info,
        }

        parse_result: Dict[str, Any] = {
            "text_content": "\n".join(text_content_parts).strip(),
            "tables": tables,
            "images": images_payload,
            "images_binary": images_binary,
            "ordered_elements": ordered_elements,
            "text_element_index_map": text_element_index_map,
            "filtered_elements_light": filtered_elements_light,
            "metadata": metadata,
            "is_converted_pdf": False,
            "converted_pdf_path": None,
        }

        logger.info(
            f"[PPTX] 解析完成: 幻灯片={len(prs.slides)}, 元素={len(ordered_elements)}, "
            f"文本={sum(1 for e in ordered_elements if e['type']=='text')}, "
            f"表格={table_count}, 图片={image_count}"
        )
        return parse_result

    def _extract_slide_content(self, slide, slide_number: int) -> Dict[str, Any]:
        """提取单张幻灯片的内容"""
        result = {
            "text": "",
            "tables": [],
            "images": [],
            "notes": "",
        }
        
        text_parts: List[str] = []
        
        # 提取所有形状的内容
        for shape in slide.shapes:
            try:
                # 跳过标题占位符（已在 _build_slide_title 中处理）
                try:
                    if slide.shapes.title and shape == slide.shapes.title:
                        continue
                except Exception:
                    # 如果无法访问 title 属性，继续处理
                    pass
                
                # 提取表格
                if shape.has_table:
                    table_data = self._extract_table_data(shape)
                    if table_data:
                        result["tables"].append(table_data)
                    continue
                
                # 提取图片
                try:
                    from pptx.enum.shapes import MSO_SHAPE_TYPE
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        image_data = self._extract_image(shape)
                        if image_data:
                            result["images"].append(image_data)
                        continue
                except Exception:
                    # 如果无法导入，尝试通过其他方式识别图片
                    try:
                        if hasattr(shape, "image") and shape.image:
                            image_data = self._extract_image(shape)
                            if image_data:
                                result["images"].append(image_data)
                            continue
                    except Exception:
                        pass
                
                # 提取文本
                text = self._extract_text_from_shape(shape)
                if text:
                    text_parts.append(text)
                    
            except Exception as e:
                logger.debug(f"[PPTX] 提取形状内容失败: {e}")
                continue
        
        result["text"] = "\n".join(text_parts)
        
        # 提取备注
        try:
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                result["notes"] = self._extract_notes(notes_slide)
        except Exception as e:
            logger.debug(f"[PPTX] 提取备注失败: {e}")
        
        return result

    def _extract_text_from_shape(self, shape) -> str:
        """从形状中提取文本"""
        try:
            if hasattr(shape, "text") and shape.text:
                return shape.text.strip()
            
            # 尝试从文本框提取
            if hasattr(shape, "text_frame"):
                text_frame = shape.text_frame
                texts = []
                for paragraph in text_frame.paragraphs:
                    para_text = paragraph.text.strip()
                    if para_text:
                        texts.append(para_text)
                return "\n".join(texts)
        except Exception as e:
            logger.debug(f"[PPTX] 提取文本失败: {e}")
        return ""

    def _extract_table_data(self, shape) -> Optional[Dict[str, Any]]:
        """提取表格数据"""
        try:
            if not shape.has_table:
                return None
            
            table = shape.table
            cells: List[List[str]] = []
            
            for row in table.rows:
                row_values: List[str] = []
                for cell in row.cells:
                    cell_text = cell.text.strip() if cell.text else ""
                    row_values.append(cell_text)
                
                # 跳过空行
                if any(val for val in row_values):
                    cells.append(row_values)
            
            if not cells:
                return None
            
            # 构建标准 table_data 结构
            table_data = {
                "cells": cells,
                "rows": len(cells),
                "columns": len(cells[0]) if cells and cells[0] else 0,
                "structure": "pptx_extracted",
                "html": None,
            }
            
            # 生成 table_text（制表符分隔）
            table_text = "\n".join("\t".join(str(cell) if cell else "" for cell in row) for row in cells)
            
            return {
                "table_data": table_data,
                "table_text": table_text,
            }
        except Exception as e:
            logger.warning(f"[PPTX] 提取表格失败: {e}")
            return None

    def _extract_image(self, shape) -> Optional[Dict[str, Any]]:
        """提取图片"""
        try:
            # 检查是否有 image 属性
            if not hasattr(shape, "image"):
                return None
            
            image = shape.image
            if not image:
                return None
            
            # 验证形状类型（如果可用）
            try:
                from pptx.enum.shapes import MSO_SHAPE_TYPE
                if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                    return None
            except Exception:
                # 如果无法导入枚举，继续处理（通过 image 属性判断）
                pass
            
            # 获取图片数据（image 已在上面获取，无需重复）
            image_bytes = image.blob
            
            if not image_bytes:
                return None
            
            # 获取图片扩展名
            image_ext = os.path.splitext(image.filename or "")[1] or ".png"
            if not image_ext.startswith("."):
                image_ext = "." + image_ext
            
            # 获取图片尺寸
            width = None
            height = None
            try:
                from PIL import Image
                img = Image.open(BytesIO(image_bytes))
                width, height = img.size
            except Exception:
                pass
            
            return {
                "data": image_bytes,
                "image_ext": image_ext,
                "width": width,
                "height": height,
            }
        except Exception as e:
            logger.warning(f"[PPTX] 提取图片失败: {e}")
            return None

    def _extract_notes(self, notes_slide) -> str:
        """提取备注页内容"""
        try:
            text_parts: List[str] = []
            for shape in notes_slide.shapes:
                text = self._extract_text_from_shape(shape)
                if text:
                    text_parts.append(text)
            return "\n".join(text_parts)
        except Exception as e:
            logger.debug(f"[PPTX] 提取备注失败: {e}")
            return ""

    def _build_slide_title(self, slide) -> str:
        """构建幻灯片标题"""
        try:
            # 优先从标题占位符提取
            if slide.shapes.title:
                title_text = slide.shapes.title.text
                if title_text and title_text.strip():
                    return title_text.strip()
            
            # 尝试从第一个文本框提取（可能是标题）
            for shape in slide.shapes:
                try:
                    from pptx.enum.shapes import MSO_SHAPE_TYPE
                    if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                        text = self._extract_text_from_shape(shape)
                        if text and len(text) < 200:  # 标题通常较短
                            return text.strip()
                except Exception:
                    # 如果无法判断类型，尝试提取文本
                    if hasattr(shape, "text_frame"):
                        text = self._extract_text_from_shape(shape)
                        if text and len(text) < 200:
                            return text.strip()
        except Exception as e:
            logger.debug(f"[PPTX] 提取标题失败: {e}")
        return ""

    def _identify_slide_layout(self, slide) -> str:
        """识别幻灯片布局"""
        try:
            if slide.slide_layout:
                return slide.slide_layout.name or "Unknown"
        except Exception:
            pass
        return "Unknown"

